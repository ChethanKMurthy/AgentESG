"""Build the Capstone presentation deck in the style of aether_esg_presentation.pptx.

Run:
    python3 build_pptx.py
Output:
    esg_intel_capstone.pptx  (repo root)
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- palette (matches reference) ----------
DARK_GREEN = RGBColor(0x1A, 0x47, 0x31)
MID_GREEN = RGBColor(0x2D, 0x6A, 0x4F)
LIGHT_GREEN = RGBColor(0x52, 0xB7, 0x88)
PALE_GREEN = RGBColor(0x74, 0xC6, 0x9D)
TEXT_DARK = RGBColor(0x1C, 0x2B, 0x2B)
MUTED = RGBColor(0x6B, 0x7F, 0x75)
LINE = RGBColor(0xD8, 0xDF, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF8, 0xFA, 0xF7)
DARK_BG = RGBColor(0x14, 0x35, 0x2B)  # richer deep green for dark slides
ACC_BLUE = RGBColor(0x24, 0x71, 0xA3)
ACC_PURPLE = RGBColor(0x7D, 0x3C, 0x98)
ACC_AMBER = RGBColor(0xF4, 0xA2, 0x61)
ACC_RED = RGBColor(0xE6, 0x39, 0x46)
ACC_VIOLET = RGBColor(0x7C, 0x3A, 0xED)
ACC_TEAL = RGBColor(0x14, 0xB8, 0xA6)
ACC_CYAN = RGBColor(0x06, 0xB6, 0xD4)

# ---------- geometry (16:9, 13.33" x 7.5") ----------
SW = 9144000
SH = 5143500
TOTAL = 17  # total slides for page-number footer


def mm_to_emu(mm: float) -> int:
    return int(mm * 36000)


def inch(x: float) -> int:
    return int(x * 914400)


def new_prs() -> Presentation:
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p


def _set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, line_rgb=None, line_w=0.5, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        # soften the corner radius
        sh.adjustments[0] = 0.12
    if fill is not None:
        _set_solid_fill(sh, fill)
    else:
        sh.fill.background()
    if line_rgb is None:
        _no_line(sh)
    else:
        _set_line(sh, line_rgb, line_w)
    sh.shadow.inherit = False
    return sh


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    font="Inter",
    size=11,
    bold=False,
    italic=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_header_strip(slide, color_bar=LIGHT_GREEN):
    # top accent bar
    bar = add_rect(slide, 0, 0, SW, Emu(55000), fill=color_bar)


def add_chip(slide, x, y, label, bg=LIGHT_GREEN, fg=WHITE):
    w = Emu(inch(1.35))
    h = Emu(inch(0.26))
    r = add_rect(slide, x, y, w, h, fill=bg, rounded=True)
    add_text(
        slide, x, y, w, h, label,
        size=7.5, bold=True, color=fg,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def add_page_number(slide, n, total=TOTAL, on_dark=False):
    add_text(
        slide, SW - inch(0.9), SH - inch(0.35), inch(0.7), inch(0.2),
        f"{n} / {total}",
        size=8, color=WHITE if on_dark else MUTED,
        align=PP_ALIGN.RIGHT,
    )


def light_slide(prs, kicker=None, n=None, bar=LIGHT_GREEN):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    bg = add_rect(s, 0, 0, SW, SH, fill=PAPER)
    # accent top bar
    add_header_strip(s, bar)
    # kicker chip
    if kicker:
        add_chip(s, inch(0.5), inch(0.22), kicker, bg=DARK_GREEN, fg=WHITE)
    if n:
        add_page_number(s, n)
    return s


def dark_slide(prs, kicker=None, n=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, SH, fill=DARK_BG)
    # decorative green blob top-right
    blob = add_rect(s, SW - inch(3.2), -inch(1.0), inch(4), inch(3), fill=MID_GREEN, rounded=True)
    blob.rotation = -12
    # soft stripe
    add_rect(s, 0, SH - Emu(55000), SW, Emu(55000), fill=LIGHT_GREEN)
    if kicker:
        add_chip(s, inch(0.5), inch(0.22), kicker, bg=WHITE, fg=DARK_GREEN)
    if n:
        add_page_number(s, n, on_dark=True)
    return s


def title_block(slide, title, sub=None, title_color=DARK_GREEN, sub_color=MUTED, top=0.7):
    add_text(
        slide, inch(0.5), inch(top), inch(12), inch(0.9),
        title, font="Georgia", size=30, bold=True, color=title_color,
    )
    if sub:
        add_text(
            slide, inch(0.5), inch(top + 0.85), inch(12), inch(0.45),
            sub, size=12.5, color=sub_color,
        )


def add_arrow(slide, x1, y1, x2, y2, color=LIGHT_GREEN, width=1.2):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    con.line.color.rgb = color
    con.line.width = Pt(width)
    # Add arrowhead via XML on line element
    lnL = con.line._get_or_add_ln()
    tail = etree.SubElement(
        lnL, qn("a:tailEnd"),
        {"type": "triangle", "w": "med", "len": "med"},
    )
    return con


def card(slide, x, y, w, h, fill=WHITE, border=LINE, rounded=True, shadow=False):
    sh = add_rect(slide, x, y, w, h, fill=fill, line_rgb=border, line_w=0.75, rounded=rounded)
    return sh


# ======================================================
# SLIDE BUILDERS
# ======================================================

def build_deck():
    prs = new_prs()

    # ---- 1. TITLE ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, SH, fill=DARK_BG)
    # decorative quarter-circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, SW - inch(5.5), -inch(2), inch(7), inch(7))
    _set_solid_fill(circ, MID_GREEN)
    _no_line(circ)
    circ2 = s.shapes.add_shape(MSO_SHAPE.OVAL, -inch(2), SH - inch(2.2), inch(4.5), inch(4.5))
    _set_solid_fill(circ2, DARK_GREEN)
    _no_line(circ2)

    add_chip(s, inch(0.6), inch(0.5), "CAPSTONE PROJECT", bg=LIGHT_GREEN, fg=WHITE)
    add_text(
        s, inch(0.6), inch(1.6), inch(12), inch(1.5),
        "Agentic ESG Intelligence",
        font="Georgia", size=58, bold=True, color=WHITE,
    )
    add_text(
        s, inch(0.6), inch(3.0), inch(12), inch(0.5),
        "An Explainable AI Platform for CSRD & ESRS Compliance",
        size=17, color=PALE_GREEN, italic=True,
    )
    add_text(
        s, inch(0.6), inch(5.1), inch(12), inch(0.4),
        "Chethan K. Murthy   •   Shashwat Saini",
        size=13, bold=True, color=WHITE,
    )
    add_text(
        s, inch(0.6), inch(5.5), inch(12), inch(0.35),
        "Guide: Dr. Prateek Verma, Associate Professor",
        size=11, color=PALE_GREEN,
    )
    add_text(
        s, inch(0.6), inch(6.6), inch(12), inch(0.35),
        "2025",
        size=12, bold=True, color=WHITE,
    )

    # ---- 2. AGENDA ----
    s = light_slide(prs, "OVERVIEW", 2)
    title_block(s, "Agenda")
    items = [
        ("01", "ESG fundamentals", "What it is and why it matters", DARK_GREEN),
        ("02", "The compliance imperative", "Regulatory landscape and academic problem", DARK_GREEN),
        ("03", "Research gap", "Why current ESG AI systems fall short", MID_GREEN),
        ("04", "Data sources", "The six datasets powering the platform", MID_GREEN),
        ("05", "System architecture", "Multi-agent orchestration, hybrid RAG, XAI", LIGHT_GREEN),
        ("06", "Demonstration & results", "Live product, benchmarks, and differentiation", LIGHT_GREEN),
    ]
    y0 = inch(2.0)
    col_w = inch(4.05)
    col_h = inch(2.0)
    for i, (n, title, desc, col) in enumerate(items):
        r, c = divmod(i, 3)
        x = inch(0.5) + c * (col_w + inch(0.1))
        y = y0 + r * (col_h + inch(0.2))
        card(s, x, y, col_w, col_h)
        add_text(s, x + inch(0.25), y + inch(0.2), inch(1.2), inch(0.6),
                 n, font="Georgia", size=22, bold=True, color=col)
        add_text(s, x + inch(0.25), y + inch(0.9), col_w - inch(0.5), inch(0.35),
                 title, size=13, bold=True, color=TEXT_DARK)
        add_text(s, x + inch(0.25), y + inch(1.3), col_w - inch(0.5), inch(0.55),
                 desc, size=10, color=MUTED)

    # ---- 3. WHAT IS ESG ----
    s = light_slide(prs, "PRIMER", 3)
    title_block(s, "What is ESG?", "A framework for evaluating a company along three non-financial dimensions.")
    pillars = [
        ("E", "Environmental", DARK_GREEN,
         "Climate impact — greenhouse-gas emissions (Scopes 1, 2, 3), energy and renewable mix, "
         "water use, biodiversity, circular economy, pollution events. Codified under ESRS E1–E5."),
        ("S", "Social", ACC_BLUE,
         "How a company treats people — own workforce, value-chain workers, affected communities, "
         "consumers. Pay equity, diversity, health & safety, human rights. ESRS S1–S4."),
        ("G", "Governance", ACC_PURPLE,
         "Quality of corporate oversight — board independence, diversity, tenure; anti-corruption, "
         "whistleblowing, business conduct, internal controls. ESRS G1."),
    ]
    y = inch(2.1)
    col_w = inch(4.1)
    for i, (letter, name, col, body) in enumerate(pillars):
        x = inch(0.5) + i * (col_w + inch(0.1))
        card(s, x, y, col_w, inch(4.2))
        add_text(s, x + inch(0.3), y + inch(0.2), inch(1.5), inch(1.6),
                 letter, font="Georgia", size=72, bold=True, color=col)
        add_text(s, x + inch(0.3), y + inch(1.7), col_w - inch(0.6), inch(0.45),
                 name, size=15, bold=True, color=col)
        add_text(s, x + inch(0.3), y + inch(2.2), col_w - inch(0.6), inch(1.9),
                 body, size=10.5, color=TEXT_DARK)

    # ---- 4. WHY COMPLIANCE IS MANDATORY ----
    s = dark_slide(prs, "MANDATE", 4)
    title_block(s, "Why ESG compliance is now mandatory",
                "A regulatory inflection point — disclosure has moved from voluntary to enforceable.",
                title_color=WHITE, sub_color=PALE_GREEN)
    metrics = [
        ("~50,000", "EU companies in CSRD scope", "Phased in from FY2024 reporting", LIGHT_GREEN),
        ("12", "ESRS standards adopted", "E1–E5, S1–S4, G1 + two cross-cutting", ACC_AMBER),
        ("Assured", "Third-party verification mandatory", "Limited assurance → reasonable", ACC_BLUE),
    ]
    y = inch(2.15)
    col_w = inch(4.1)
    for i, (big, cap, sub, col) in enumerate(metrics):
        x = inch(0.5) + i * (col_w + inch(0.1))
        card(s, x, y, col_w, inch(3.0), fill=RGBColor(0x0E, 0x24, 0x1C), border=MID_GREEN)
        add_text(s, x + inch(0.3), y + inch(0.4), col_w - inch(0.6), inch(1.2),
                 big, font="Georgia", size=52, bold=True, color=col)
        add_text(s, x + inch(0.3), y + inch(1.7), col_w - inch(0.6), inch(0.5),
                 cap, size=13, bold=True, color=WHITE)
        add_text(s, x + inch(0.3), y + inch(2.3), col_w - inch(0.6), inch(0.5),
                 sub, size=10, color=PALE_GREEN)

    # regulatory timeline
    tl_y = inch(5.5)
    add_text(s, inch(0.5), tl_y - inch(0.4), inch(12), inch(0.3),
             "Regulatory convergence across jurisdictions",
             size=11, bold=True, color=PALE_GREEN)
    milestones = [
        ("2015", "Paris Agreement"),
        ("2017", "TCFD"),
        ("2022", "CSRD adopted"),
        ("2023", "ISSB S1/S2"),
        ("2024", "CSRD Y1 reporting"),
        ("2026", "Wave-2 expansion"),
    ]
    ms_w = inch(12) / len(milestones)
    # timeline line
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  inch(0.5), tl_y + inch(0.18),
                                  inch(12.5), tl_y + inch(0.18))
    line.line.color.rgb = LIGHT_GREEN
    line.line.width = Pt(1.2)
    for i, (yr, ev) in enumerate(milestones):
        cx = inch(0.5) + i * ms_w + ms_w / 2
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(inch(0.07)), tl_y + inch(0.11), inch(0.14), inch(0.14))
        _set_solid_fill(dot, WHITE)
        _no_line(dot)
        add_text(s, cx - Emu(inch(0.8)), tl_y - inch(0.1), inch(1.6), inch(0.25),
                 yr, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cx - Emu(inch(0.9)), tl_y + inch(0.36), inch(1.8), inch(0.35),
                 ev, size=9, color=PALE_GREEN, align=PP_ALIGN.CENTER)

    # ---- 5. PROBLEM STATEMENT ----
    s = dark_slide(prs, "PROBLEM", 5)
    title_block(s, "The academic problem",
                "Why compliance at scale is fundamentally hard — beyond the investor story.",
                title_color=WHITE, sub_color=PALE_GREEN)
    probs = [
        ("Unstructured disclosure",
         "Annual reports range from 40 to 400+ pages of narrative, tables, and footnotes — "
         "no canonical schema, heavy NLP burden."),
        ("Sector heterogeneity",
         "A steel mill and a mobile operator disclose different metrics, against different "
         "thresholds, under the same directive."),
        ("Double materiality",
         "ESRS requires evaluating impact on company AND by company — a non-trivial reasoning "
         "problem that generic LLMs cannot do reliably."),
        ("Audit exposure",
         "Reports must carry limited/reasonable assurance. Every figure must be traceable "
         "back to source evidence — AI hallucinations are not tolerated."),
    ]
    y = inch(2.1)
    col_w = inch(6.1)
    for i, (t, body) in enumerate(probs):
        r, c = divmod(i, 2)
        x = inch(0.5) + c * (col_w + inch(0.15))
        yy = y + r * (inch(1.9))
        card(s, x, yy, col_w, inch(1.65), fill=RGBColor(0x0E, 0x24, 0x1C), border=MID_GREEN)
        add_text(s, x + inch(0.3), yy + inch(0.2), col_w - inch(0.6), inch(0.4),
                 t, size=14, bold=True, color=ACC_AMBER)
        add_text(s, x + inch(0.3), yy + inch(0.7), col_w - inch(0.6), inch(0.95),
                 body, size=10.5, color=WHITE)

    # ---- 6. RESEARCH GAP ----
    s = light_slide(prs, "RESEARCH GAP", 6)
    title_block(s, "What current ESG-AI systems miss",
                "Existing tools either score silently or generate ungrounded prose. Neither passes audit.")
    table_rows = [
        ["Capability", "Legacy ESG SaaS", "Generic LLM wrappers", "Our system"],
        ["Deterministic rule evaluation", "Partial", "None", "Yes (10+ versioned rules)"],
        ["Retrieval-grounded reasoning", "None", "Partial", "Hybrid (FAISS + BM25 + RRF)"],
        ["Verifiable citations", "None", "Inconsistent", "Every claim carries [n]"],
        ["Structured XAI for roadmap", "None", "None", "Drivers + factors + narrative"],
        ["Multi-agent orchestration", "None", "Ad-hoc", "Typed context + audit trace"],
        ["Industry-scoped benchmarking", "Aggregated", "None", "Per-metric percentile + median"],
    ]
    left = inch(0.5)
    top = inch(2.1)
    total_w = inch(12.3)
    row_h = inch(0.48)
    col_ws = [inch(3.2), inch(2.6), inch(2.6), inch(3.9)]
    for r_idx, row in enumerate(table_rows):
        x = left
        for c_idx, cell in enumerate(row):
            is_header = r_idx == 0
            is_us = c_idx == 3 and not is_header
            fill = DARK_GREEN if is_header else (PALE_GREEN if is_us else WHITE)
            border = DARK_GREEN if is_header else LINE
            cw = col_ws[c_idx]
            card(s, x, top + r_idx * row_h, cw, row_h, fill=fill, border=border, rounded=False)
            add_text(
                s, x + inch(0.15), top + r_idx * row_h, cw - inch(0.3), row_h,
                cell,
                size=10,
                bold=is_header or is_us,
                color=WHITE if is_header else (DARK_GREEN if is_us else TEXT_DARK),
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += cw

    # ---- 7. DATA SOURCES ----
    s = light_slide(prs, "DATA", 7)
    title_block(s, "Data sources",
                "Six dataset classes drive the platform — per the project's dataset specification.")
    rows = [
        ("Company master", "1,000 EU-focused entities",
         "company_id, sector, industry, revenue, employees, listing, size_band — identity and peer-grouping keys.",
         DARK_GREEN),
        ("ESG metrics", "Keyed by (company_id, year)",
         "Scopes 1/2/3, energy, renewable %, water, waste, female ratio, board independence, incidents, ESG score.",
         MID_GREEN),
        ("Regulatory corpus (RAG)", "CSRD, ESRS, GHG Protocol, GRI, SASB, ISSB",
         "PDF and markdown ingested, chunked at ~800 tokens, embedded, and indexed for hybrid retrieval.",
         LIGHT_GREEN),
        ("Regulatory rules", "10 declarative, versioned rules",
         "CSRD applicability (EU large, listed, non-EU), ESRS E1/E3/E5/G1/S1 thresholds with citations.",
         ACC_BLUE),
        ("Industry benchmarks", "Per-sector / per-industry / per-country",
         "Derived on demand from the metrics table — p25/p50/p75, direction-aware percentile ranks.",
         ACC_PURPLE),
        ("Audit & interaction logs", "Request-ID correlated",
         "Every agent step, rule firing, retrieval call, and LLM call is written to Postgres with latency.",
         ACC_AMBER),
    ]
    y0 = inch(2.0)
    col_w = inch(4.1)
    row_h = inch(1.55)
    for i, (title, sub, body, col) in enumerate(rows):
        r, c = divmod(i, 3)
        x = inch(0.5) + c * (col_w + inch(0.1))
        yy = y0 + r * (row_h + inch(0.2))
        card(s, x, yy, col_w, row_h)
        # colored bar on left
        bar = add_rect(s, x, yy, Emu(inch(0.12)), row_h, fill=col, rounded=False)
        add_text(s, x + inch(0.3), yy + inch(0.15), col_w - inch(0.5), inch(0.35),
                 title, size=12, bold=True, color=col)
        add_text(s, x + inch(0.3), yy + inch(0.55), col_w - inch(0.5), inch(0.3),
                 sub, size=8.5, color=MUTED)
        add_text(s, x + inch(0.3), yy + inch(0.85), col_w - inch(0.5), inch(0.7),
                 body, size=9.5, color=TEXT_DARK)

    # ---- 8. SYSTEM ARCHITECTURE ----
    s = light_slide(prs, "ARCHITECTURE", 8)
    title_block(s, "System architecture",
                "A layered, stateless platform. Every arrow is auditable.")

    # Draw a layered block diagram
    # Layers: Client → API / Middleware → Agents → RAG + Rule Engine + Analytics → Data
    def layer(x, y, w, h, title, boxes, tone=LIGHT_GREEN):
        # title strip
        add_text(s, x, y, w, inch(0.3), title,
                 size=9.5, bold=True, color=tone, align=PP_ALIGN.CENTER)
        card(s, x, y + inch(0.3), w, h - inch(0.3), border=tone, rounded=True)
        bx = x + inch(0.15)
        bw = (w - inch(0.2) - inch(0.1) * (len(boxes) - 1)) / len(boxes)
        for i, (lbl, sub) in enumerate(boxes):
            bx_i = bx + i * (bw + inch(0.1))
            inner = card(s, bx_i, y + inch(0.45), bw, h - inch(0.55),
                         fill=PAPER, border=LINE)
            add_text(s, bx_i, y + inch(0.45), bw, inch(0.35),
                     lbl, size=10, bold=True, color=TEXT_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, bx_i + inch(0.08), y + inch(0.85), bw - inch(0.16), inch(0.55),
                     sub, size=8.2, color=MUTED, align=PP_ALIGN.CENTER)

    x0 = inch(0.5)
    w0 = inch(12.3)
    top = inch(1.9)
    lh = inch(1.0)
    gap = inch(0.2)

    layer(x0, top, w0, lh, "CLIENT",
          [("React SPA", "home, upload, dashboard"),
           ("Verdant chat widget", "streaming + briefing"),
           ("nginx proxy", "/ and /api/*")],
          tone=LIGHT_GREEN)

    layer(x0, top + lh + gap, w0, lh, "API & MIDDLEWARE",
          [("FastAPI", "async routes"),
           ("Audit middleware", "request-id, latency"),
           ("Pydantic schemas", "contract validation"),
           ("CORS / auth", "origin allow-list")],
          tone=MID_GREEN)

    layer(x0, top + 2 * (lh + gap), w0, lh, "AGENT ORCHESTRATION",
          [("Data", "company + metrics"),
           ("Retrieval", "hybrid RAG"),
           ("Analysis", "rule engine"),
           ("Report", "Claude prompt"),
           ("Audit", "trace rows")],
          tone=DARK_GREEN)

    layer(x0, top + 3 * (lh + gap), w0, lh, "KNOWLEDGE & DATA",
          [("Postgres", "companies, metrics, audit"),
           ("FAISS + BM25", "vector + sparse"),
           ("HF embed model", "MiniLM-L6-v2"),
           ("Anthropic Claude", "Opus 4.7")],
          tone=ACC_BLUE)

    # ---- 9. MULTI-AGENT WORKFLOW ----
    s = light_slide(prs, "AGENTS", 9)
    title_block(s, "Multi-agent workflow",
                "Five agents, shared typed context, auto-traced. Order is deterministic.")
    # flow of 5 agents with arrows
    agents = [
        ("Data", "Loads company +\nlatest metrics\nfrom Postgres", DARK_GREEN),
        ("Retrieval", "Queries hybrid\nRAG for\nregulatory context", MID_GREEN),
        ("Analysis", "Runs 10+ rule\nengine checks\ndeterministically", LIGHT_GREEN),
        ("Report", "Grounded Claude\nsynthesis with\ninline [n] cites", ACC_BLUE),
        ("Audit", "Writes 6 audit\nrows with\nrequest-id", ACC_PURPLE),
    ]
    n = len(agents)
    total_w = inch(12)
    slot = total_w / n
    y = inch(2.3)
    box_w = slot - inch(0.2)
    box_h = inch(2.2)
    for i, (name, body, col) in enumerate(agents):
        x = inch(0.65) + i * slot
        # colored top stripe
        card(s, x, y, box_w, box_h)
        stripe = add_rect(s, x, y, box_w, inch(0.35), fill=col, rounded=False)
        add_text(s, x, y, box_w, inch(0.35), name.upper(),
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + inch(0.15), y + inch(0.5), box_w - inch(0.3), inch(1.6),
                 body, size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        # arrow to next
        if i < n - 1:
            ax1 = x + box_w + Emu(inch(0.02))
            ax2 = x + slot - inch(0.02)
            cy = y + box_h / 2
            add_arrow(s, ax1, cy, ax2, cy, color=MUTED)

    # shared AgentContext strip
    ctx_y = y + box_h + inch(0.4)
    card(s, inch(0.65), ctx_y, total_w, inch(1.2), fill=RGBColor(0xEC, 0xF3, 0xEE), border=LIGHT_GREEN)
    add_text(s, inch(0.85), ctx_y + inch(0.15), total_w - inch(0.4), inch(0.35),
             "Shared AgentContext", size=11, bold=True, color=DARK_GREEN)
    add_text(s, inch(0.85), ctx_y + inch(0.55), total_w - inch(0.4), inch(0.55),
             "company_id · company · latest_metrics · retrieved[] · rule_results[] · rule_summary · report · trace[] · errors[] · request_id",
             size=9.5, color=TEXT_DARK, font="Consolas")

    # ---- 10. HYBRID RAG ----
    s = dark_slide(prs, "RETRIEVAL", 10)
    title_block(s, "Hybrid retrieval-augmented generation",
                "Dense + sparse signals fused with RRF, diversified with MMR.",
                title_color=WHITE, sub_color=PALE_GREEN)

    # Left: pipeline (5 stages vertical)
    stages = [
        ("Query", "user or agent question"),
        ("Dense (FAISS)", "MiniLM-L6-v2, 384-d, IP"),
        ("Sparse (BM25)", "token overlap over same chunks"),
        ("Reciprocal-Rank Fusion", "weights: 0.6 dense, 0.4 sparse, k=60"),
        ("MMR reranker", "λ=0.7 for relevance/diversity"),
        ("Top-K chunks → LLM", "with inline [n] citations"),
    ]
    x = inch(0.7)
    y = inch(2.15)
    w = inch(5.3)
    h = inch(0.62)
    for i, (t, sub) in enumerate(stages):
        yy = y + i * (h + inch(0.08))
        card(s, x, yy, w, h, fill=RGBColor(0x0E, 0x24, 0x1C), border=MID_GREEN)
        add_text(s, x + inch(0.2), yy, inch(2.6), h,
                 t, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + inch(2.9), yy, w - inch(3.1), h,
                 sub, size=10, color=PALE_GREEN, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            add_arrow(s, x + w / 2, yy + h, x + w / 2, yy + h + inch(0.08),
                      color=LIGHT_GREEN, width=1.5)

    # Right: formula + rationale
    right_x = inch(6.4)
    right_y = inch(2.15)
    right_w = inch(6.3)
    right_h = inch(4.2)
    card(s, right_x, right_y, right_w, right_h, fill=RGBColor(0x0E, 0x24, 0x1C), border=MID_GREEN)
    add_text(s, right_x + inch(0.3), right_y + inch(0.25), right_w - inch(0.6), inch(0.4),
             "Why hybrid fusion?", size=14, bold=True, color=LIGHT_GREEN)
    add_text(s, right_x + inch(0.3), right_y + inch(0.75), right_w - inch(0.6), inch(0.95),
             "Dense embeddings capture semantic similarity (e.g., 'carbon intensity' ≈ "
             "'emissions per unit revenue'). BM25 captures exact lexical matches on "
             "regulation numbers and specific terms ('Scope 3', 'Article 19a'). "
             "Regulatory QA needs both.",
             size=10.5, color=WHITE)
    add_text(s, right_x + inch(0.3), right_y + inch(1.85), right_w - inch(0.6), inch(0.35),
             "RRF score", size=12, bold=True, color=LIGHT_GREEN)
    add_text(s, right_x + inch(0.3), right_y + inch(2.25), right_w - inch(0.6), inch(0.5),
             "score(doc) = w_d / (k + rank_d(doc))  +  w_s / (k + rank_s(doc))",
             size=11, color=PALE_GREEN, font="Consolas")
    add_text(s, right_x + inch(0.3), right_y + inch(2.95), right_w - inch(0.6), inch(0.35),
             "MMR rerank", size=12, bold=True, color=LIGHT_GREEN)
    add_text(s, right_x + inch(0.3), right_y + inch(3.35), right_w - inch(0.6), inch(0.7),
             "MMR(d) = λ · rel(d,q) − (1 − λ) · maxₛ sim(d, s)  →  favours relevant passages "
             "that add new information rather than restating prior hits.",
             size=10.5, color=WHITE)

    # ---- 11. RULE ENGINE + XAI ----
    s = light_slide(prs, "XAI", 11)
    title_block(s, "Rule engine & roadmap explainability",
                "Deterministic decisions on the left. Transparent explanations for every recommended action.")
    # Left card: rule data model
    lx = inch(0.5)
    ly = inch(2.0)
    lw = inch(6.0)
    lh = inch(4.6)
    card(s, lx, ly, lw, lh, fill=RGBColor(0xEC, 0xF3, 0xEE), border=LIGHT_GREEN)
    add_text(s, lx + inch(0.3), ly + inch(0.25), lw - inch(0.6), inch(0.4),
             "Rule as data (not code)", size=14, bold=True, color=DARK_GREEN)
    code = (
        "Rule(\n"
        "  id='CSRD_APPLICABILITY_LARGE_EU',\n"
        "  category='applicability',\n"
        "  severity='info',\n"
        "  when=lambda c, m: is_eu(c) and meets(c) >= 2,\n"
        "  message=lambda c, m: f'CSRD applies: {meets(c)}/3 thresholds met',\n"
        "  citations=['CSRD Art. 3; Directive (EU) 2022/2464']\n"
        ")"
    )
    add_text(s, lx + inch(0.3), ly + inch(0.8), lw - inch(0.6), inch(3.6),
             code, size=10, color=TEXT_DARK, font="Consolas", line_spacing=1.3)

    # Right: XAI card
    rx = inch(6.8)
    ry = inch(2.0)
    rw = inch(6.0)
    card(s, rx, ry, rw, lh, fill=WHITE, border=LINE)
    add_text(s, rx + inch(0.3), ry + inch(0.25), rw - inch(0.6), inch(0.4),
             "Roadmap explanation (per item)", size=14, bold=True, color=DARK_GREEN)
    # inner panels
    panels = [
        ("DRIVERS",
         "rule: GOV_BOARD_INDEPENDENCE_LOW  sev=high  current=36.26  target=50\n"
         "benchmark: board_independence_percent  current=36.26  target=61.1 (sector median)"),
        ("PRIORITY FACTORS",
         "severity_rank=4   gap_count=2   rule_drivers=1   benchmark_drivers=1   weight_sum=8.0"),
        ("HORIZON RATIONALE",
         "Horizon 0–3m assigned because the highest-severity driver is 'high'."),
        ("VERDANT NARRATIVE  (LLM)",
         "Board independence at 36.26% is materially below the 50% threshold and in the "
         "bottom 2nd percentile versus peers (median 61.1%, n=94). Under ESRS G1 this "
         "represents a disclosure weakness requiring near-term remediation."),
    ]
    yy = ry + inch(0.85)
    for (label, body) in panels:
        card(s, rx + inch(0.3), yy, rw - inch(0.6), inch(0.9),
             fill=PAPER, border=LINE)
        add_text(s, rx + inch(0.45), yy + inch(0.05), rw - inch(0.9), inch(0.3),
                 label, size=8.5, bold=True, color=MID_GREEN)
        add_text(s, rx + inch(0.45), yy + inch(0.35), rw - inch(0.9), inch(0.55),
                 body, size=9, color=TEXT_DARK, font="Consolas", line_spacing=1.25)
        yy += inch(0.93)

    # ---- 12. DISCLOSURE EXTRACTION + BENCHMARKING ----
    s = light_slide(prs, "INGESTION", 12)
    title_block(s, "PDF disclosure extraction & industry benchmarking",
                "Upload an annual report; the system extracts KPIs and ranks the company against industry peers.")
    # Left flow: PDF → text → LLM → structured → DB
    flow = [("PDF", "annual / sustainability\nreport"),
            ("Text extract", "pypdf, 120K char cap\nhead+tail trim"),
            ("LLM extract", "Claude with strict\nJSON schema"),
            ("Sanitise", "type coercion, unit\nvalidation"),
            ("Upsert", "Postgres company\n+ metric row")]
    lx = inch(0.5)
    ly = inch(2.1)
    lw = inch(6.0)
    col_w = (lw - inch(0.5)) / len(flow)
    for i, (t, sub) in enumerate(flow):
        x = lx + i * col_w
        card(s, x + inch(0.05), ly, col_w - inch(0.1), inch(1.7))
        add_text(s, x + inch(0.05), ly + inch(0.25), col_w - inch(0.1), inch(0.4),
                 t, size=11, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
        add_text(s, x + inch(0.15), ly + inch(0.8), col_w - inch(0.3), inch(0.85),
                 sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            cy = ly + inch(0.85)
            add_arrow(s, x + col_w - inch(0.08), cy, x + col_w + inch(0.06), cy, color=LIGHT_GREEN)

    # Right: benchmark bar viz
    rx = inch(0.5)
    ry = inch(4.1)
    rw = inch(12.3)
    card(s, rx, ry, rw, inch(2.6))
    add_text(s, rx + inch(0.3), ry + inch(0.2), rw - inch(0.6), inch(0.35),
             "Direction-aware percentile vs industry peers",
             size=12, bold=True, color=DARK_GREEN)
    add_text(s, rx + inch(0.3), ry + inch(0.55), rw - inch(0.6), inch(0.3),
             "Higher-is-better metrics (renewable %, board indep.) vs lower-is-better metrics (scope emissions, incidents).",
             size=9.5, color=MUTED)
    rows = [
        ("Renewable energy", 0.86, "leader",  LIGHT_GREEN),
        ("Board independence", 0.02, "laggard", ACC_RED),
        ("Scope 3 emissions",   0.46, "peer",    ACC_AMBER),
        ("Waste recycled",      0.89, "leader",  LIGHT_GREEN),
    ]
    bar_y0 = ry + inch(1.05)
    row_h = inch(0.35)
    for i, (name, pct, status, col) in enumerate(rows):
        yy = bar_y0 + i * row_h
        add_text(s, rx + inch(0.3), yy, inch(2.2), row_h - inch(0.05),
                 name, size=10, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        track = add_rect(s, rx + inch(2.6), yy + inch(0.1), inch(8), inch(0.16),
                         fill=RGBColor(0xEA, 0xEE, 0xE8), rounded=True)
        fill_w = int(inch(8) * pct)
        fill_bar = add_rect(s, rx + inch(2.6), yy + inch(0.1), fill_w, inch(0.16),
                            fill=col, rounded=True)
        # median tick
        median_x = rx + inch(2.6) + int(inch(8) * 0.5)
        tick = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      median_x, yy + inch(0.05),
                                      median_x, yy + inch(0.31))
        tick.line.color.rgb = MUTED
        tick.line.width = Pt(0.5)
        add_text(s, rx + inch(10.8), yy, inch(1.6), row_h - inch(0.05),
                 status, size=9, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)

    # ---- 13. PRODUCT DEMO (mock UI) ----
    s = dark_slide(prs, "DEMO", 13)
    title_block(s, "Product demonstration",
                "Three surfaces, one coherent workstation.",
                title_color=WHITE, sub_color=PALE_GREEN)

    # three mock screens
    demos = [
        ("Upload PDF",
         ["DROP YOUR PDF HERE",
          "(dumbledore_industries_fy2024.pdf)",
          "",
          "1 • Parsing PDF",
          "2 • Reading disclosure",
          "3 • Extracting KPIs",
          "4 • Persisting company",
          "5 • Benchmarking vs industry"]),
        ("Company dashboard",
         ["Dumbledore Industries plc",
          "UK · Energy/Utilities · Renewable Energy",
          "",
          "ESG score:  78.9",
          "Renewable:  86.4%",
          "Board indep: 72%",
          "Peers:   8 (industry)",
          "Rank:    leader"]),
        ("Agent Verdant",
         ["Codename: Verdant",
          "corpus: 4 docs · 8 chunks",
          "",
          "You: Does CSRD apply here?",
          "",
          "Verdant: Yes. The entity meets",
          "employees + revenue thresholds",
          "per CSRD Art. 3 [1]. Large-",
          "undertaking path applies."]),
    ]
    y = inch(2.2)
    w = inch(4.1)
    h = inch(4.3)
    for i, (title, lines) in enumerate(demos):
        x = inch(0.5) + i * (w + inch(0.15))
        # device frame
        card(s, x, y, w, h, fill=RGBColor(0x0E, 0x24, 0x1C), border=LIGHT_GREEN)
        # top title bar
        add_rect(s, x, y, w, inch(0.3), fill=LIGHT_GREEN, rounded=False)
        add_text(s, x, y, w, inch(0.3), title.upper(),
                 size=9.5, bold=True, color=DARK_GREEN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # screen body
        body_x = x + inch(0.2)
        body_y = y + inch(0.45)
        body_w = w - inch(0.4)
        body_h = h - inch(0.6)
        card(s, body_x, body_y, body_w, body_h, fill=PAPER, border=LINE)
        inner_y = body_y + inch(0.2)
        for ln in lines:
            font_size = 10 if ln and not ln.startswith(("1 ", "2 ", "3 ", "4 ", "5 ")) else 9
            bold = ln == lines[0] or ln.startswith("You:") or ln.startswith("Verdant:")
            col = DARK_GREEN if bold else TEXT_DARK
            add_text(s, body_x + inch(0.15), inner_y, body_w - inch(0.3), inch(0.3),
                     ln, size=font_size, bold=bold, color=col, font="Consolas")
            inner_y += inch(0.3)

    # ---- 14. TECH STACK ----
    s = light_slide(prs, "STACK", 14)
    title_block(s, "Technology & AI/ML stack",
                "Everything runs locally via Docker Compose; horizontally scalable.")
    cols = [
        ("BACKEND", DARK_GREEN, [
            "Python 3.11, FastAPI (async)",
            "SQLAlchemy 2.0 + asyncpg",
            "Pydantic v2 / Pydantic-Settings",
            "structlog (JSON logs)",
            "Alembic migrations",
            "Gunicorn + UvicornWorker (prod)",
        ]),
        ("AI / ML", MID_GREEN, [
            "Anthropic Claude Opus 4.7",
            "sentence-transformers MiniLM-L6-v2",
            "FAISS-CPU (inner product)",
            "rank-bm25 sparse retrieval",
            "Reciprocal-Rank Fusion + MMR",
            "pypdf + python-docx parsing",
        ]),
        ("FRONTEND", LIGHT_GREEN, [
            "React 18 + TypeScript + Vite",
            "Tailwind CSS (custom tokens)",
            "Framer Motion animations",
            "Recharts visualisation",
            "TanStack React Query",
            "nginx static + reverse proxy",
        ]),
        ("INFRA", ACC_BLUE, [
            "Docker Compose (multi-stage)",
            "Postgres 16-alpine (pg_isready)",
            "Docker volumes for model cache",
            "Healthchecks + restart policies",
            "Cloudflare Tunnel (demo access)",
            "AWS ECS/RDS/S3 templates ready",
        ]),
    ]
    y0 = inch(1.95)
    col_w = inch(3.05)
    col_h = inch(4.55)
    for i, (title, col, items) in enumerate(cols):
        x = inch(0.5) + i * (col_w + inch(0.1))
        card(s, x, y0, col_w, col_h)
        add_rect(s, x, y0, col_w, inch(0.4), fill=col, rounded=False)
        add_text(s, x, y0, col_w, inch(0.4), title,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ty = y0 + inch(0.6)
        for it in items:
            # bullet
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + inch(0.25), ty + inch(0.09), inch(0.08), inch(0.08))
            _set_solid_fill(dot, col)
            _no_line(dot)
            add_text(s, x + inch(0.45), ty, col_w - inch(0.6), inch(0.55),
                     it, size=10, color=TEXT_DARK)
            ty += inch(0.52)

    # ---- 15. RESULTS & DIFFERENTIATION ----
    s = dark_slide(prs, "RESULTS", 15)
    title_block(s, "Results & differentiation",
                "Measured outcomes on the included synthetic + real peer universe.",
                title_color=WHITE, sub_color=PALE_GREEN)
    kpi = [
        ("1,000",  "companies benchmarked across 10 sectors, 80+ industries", LIGHT_GREEN),
        ("10",     "CSRD / ESRS rules evaluated per company (versioned)", ACC_AMBER),
        ("14",     "metrics extracted from PDF disclosures in zero-shot mode", ACC_BLUE),
        ("100%",   "LLM answers carry inline [n] citations back to source", PALE_GREEN),
    ]
    y = inch(2.1)
    col_w = inch(3.07)
    for i, (big, cap, col) in enumerate(kpi):
        x = inch(0.5) + i * (col_w + inch(0.1))
        card(s, x, y, col_w, inch(2.0), fill=RGBColor(0x0E, 0x24, 0x1C), border=MID_GREEN)
        add_text(s, x + inch(0.25), y + inch(0.25), col_w - inch(0.5), inch(0.9),
                 big, font="Georgia", size=42, bold=True, color=col)
        add_text(s, x + inch(0.25), y + inch(1.2), col_w - inch(0.5), inch(0.75),
                 cap, size=10, color=WHITE)

    # novelty call-outs
    call_y = inch(4.35)
    call_w = (inch(12.3) - inch(0.2)) / 3
    callouts = [
        ("Hybrid RAG with RRF",
         "Not a single-vector store; verifiable dense + sparse fusion."),
        ("Typed multi-agent orchestration",
         "AgentContext + auto-trace, not a free-form LangChain graph."),
        ("Structured XAI on roadmap",
         "Every recommended action carries drivers, factors, and narrative."),
    ]
    for i, (t, b) in enumerate(callouts):
        x = inch(0.5) + i * (call_w + inch(0.1))
        card(s, x, call_y, call_w, inch(1.75), fill=RGBColor(0x0E, 0x24, 0x1C), border=LIGHT_GREEN)
        add_text(s, x + inch(0.25), call_y + inch(0.2), call_w - inch(0.5), inch(0.4),
                 t, size=12, bold=True, color=LIGHT_GREEN)
        add_text(s, x + inch(0.25), call_y + inch(0.7), call_w - inch(0.5), inch(1),
                 b, size=10, color=WHITE)

    # ---- 16. CONCLUSION ----
    s = light_slide(prs, "CONCLUSION", 16)
    title_block(s, "Conclusion & future work",
                "A research prototype that demonstrates auditable agentic AI for regulatory reasoning.")
    contribs = [
        ("Conceptual",
         "Formalised ESG compliance reasoning as a typed multi-agent pipeline over structured + "
         "unstructured evidence.",
         DARK_GREEN),
        ("Technical",
         "Hybrid RAG with direction-aware benchmarking and a data-driven rule engine, all under "
         "a single audit fabric.",
         MID_GREEN),
        ("Applied",
         "A working platform with PDF-driven extraction, industry-scoped peer ranking, and "
         "explainable roadmap generation.",
         LIGHT_GREEN),
    ]
    y = inch(2.0)
    col_w = inch(4.05)
    for i, (t, b, col) in enumerate(contribs):
        x = inch(0.5) + i * (col_w + inch(0.15))
        card(s, x, y, col_w, inch(2.3))
        add_rect(s, x, y, col_w, inch(0.35), fill=col, rounded=False)
        add_text(s, x, y, col_w, inch(0.35), t.upper(),
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + inch(0.25), y + inch(0.55), col_w - inch(0.5), inch(1.6),
                 b, size=11, color=TEXT_DARK)

    # Future work
    fw_y = inch(4.6)
    card(s, inch(0.5), fw_y, inch(12.3), inch(2.2), fill=RGBColor(0xEC, 0xF3, 0xEE), border=LIGHT_GREEN)
    add_text(s, inch(0.75), fw_y + inch(0.2), inch(12), inch(0.4),
             "Future work", size=14, bold=True, color=DARK_GREEN)
    fw_items = [
        "Pluggable rule authoring via YAML DSL with versioning and unit tests.",
        "Multi-year trajectory modelling and scenario analysis against 1.5 °C / 2 °C pathways.",
        "Image + table extraction for non-textual disclosure artefacts (OCR + layout).",
        "SOC 2 / ISO 27001-aligned audit export (signed trace, tamper-evident).",
        "Evaluation against a held-out set of real CSRD reports with human-graded findings.",
    ]
    for i, it in enumerate(fw_items):
        r, c = divmod(i, 2)
        x = inch(0.85) + c * inch(6.0)
        yy = fw_y + inch(0.7) + r * inch(0.35)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + inch(0.1), inch(0.08), inch(0.08))
        _set_solid_fill(dot, MID_GREEN)
        _no_line(dot)
        add_text(s, x + inch(0.2), yy, inch(5.7), inch(0.35),
                 it, size=10.5, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # ---- 17. THANK YOU ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, SH, fill=DARK_BG)
    # decorative
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, -inch(2), -inch(2), inch(5), inch(5))
    _set_solid_fill(circ, MID_GREEN)
    _no_line(circ)
    circ2 = s.shapes.add_shape(MSO_SHAPE.OVAL, SW - inch(3), SH - inch(3), inch(5), inch(5))
    _set_solid_fill(circ2, DARK_GREEN)
    _no_line(circ2)
    add_text(s, inch(0.5), inch(2.0), inch(12), inch(1.5),
             "Thank You",
             font="Georgia", size=72, bold=True, color=WHITE)
    add_text(s, inch(0.5), inch(3.7), inch(12), inch(0.5),
             "Agentic ESG Intelligence — auditable AI for regulatory reasoning.",
             size=15, italic=True, color=PALE_GREEN)
    # SDGs
    sdg_y = inch(4.7)
    add_text(s, inch(0.5), sdg_y, inch(12), inch(0.3),
             "Aligned with UN Sustainable Development Goals:",
             size=11, color=WHITE)
    sdgs = [("SDG 12", "Responsible Consumption"),
            ("SDG 13", "Climate Action"),
            ("SDG 16", "Peace, Justice, Institutions")]
    for i, (n, nm) in enumerate(sdgs):
        x = inch(1.0) + i * inch(4.2)
        card(s, x, sdg_y + inch(0.4), inch(3.5), inch(0.9),
             fill=RGBColor(0x0E, 0x24, 0x1C), border=LIGHT_GREEN)
        add_text(s, x + inch(0.25), sdg_y + inch(0.5), inch(3), inch(0.35),
                 n, size=14, bold=True, color=LIGHT_GREEN)
        add_text(s, x + inch(0.25), sdg_y + inch(0.85), inch(3), inch(0.35),
                 nm, size=10, color=WHITE)

    add_text(s, inch(0.5), SH - inch(0.7), inch(12), inch(0.4),
             "Chethan K. Murthy   •   Shashwat Saini   •   Guide: Dr. Prateek Verma",
             size=11, bold=True, color=WHITE)

    # save
    out = "esg_intel_capstone.pptx"
    prs.save(out)
    print("wrote", out, "with", len(prs.slides), "slides")


if __name__ == "__main__":
    build_deck()
